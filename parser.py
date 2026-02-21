"""
parser.py
---------
Loads documents from ALL supported source types into LangChain Documents.

Supported formats:
  Unstructured  → .txt, .md, .pdf, .docx, .pptx
  Semi-structured → .json, .xml
  Structured    → .csv, .sqlite / .db

Returns List[Document] ready for chunker.py.
"""

import os
import csv
import json
import sqlite3
import xml.etree.ElementTree as ET
from typing import List

from langchain_core.documents import Document
from langchain_community.document_loaders import (
    TextLoader,
    DirectoryLoader,
    PyMuPDFLoader,
    PyPDFLoader,
    Docx2txtLoader,
)


# ─────────────────────────────────────────────────────────────────────────────
# Unstructured
# ─────────────────────────────────────────────────────────────────────────────

def load_text_file(file_path: str) -> List[Document]:
    loader = TextLoader(file_path, encoding="utf-8")
    return loader.load()


def load_text_directory(directory_path: str, glob: str = "**/*.txt") -> List[Document]:
    loader = DirectoryLoader(
        directory_path,
        glob=glob,
        loader_cls=TextLoader,
        loader_kwargs={"encoding": "utf-8"},
        show_progress=True,
    )
    return loader.load()


def load_pdf(file_path: str, method: str = "pymupdf") -> List[Document]:
    loader = PyMuPDFLoader(file_path) if method == "pymupdf" else PyPDFLoader(file_path)
    return loader.load()


def load_docx(file_path: str) -> List[Document]:
    loader = Docx2txtLoader(file_path)
    docs = loader.load()
    for doc in docs:
        doc.metadata.update({"source": file_path, "data_type": "docx"})
    return docs


def load_pptx(file_path: str) -> List[Document]:
    from pptx import Presentation as PptxPresentation
    prs = PptxPresentation(file_path)
    docs = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            docs.append(Document(
                page_content="\n".join(texts),
                metadata={"source": file_path, "slide": i, "data_type": "pptx"},
            ))
    return docs


# ─────────────────────────────────────────────────────────────────────────────
# Semi-structured
# ─────────────────────────────────────────────────────────────────────────────

def load_json(file_path: str) -> List[Document]:
    with open(file_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    def flatten(obj, prefix="") -> List[str]:
        lines = []
        if isinstance(obj, dict):
            for k, v in obj.items():
                lines.extend(flatten(v, f"{prefix}{k}: " if not prefix else f"{prefix}.{k}: "))
        elif isinstance(obj, list):
            for i, item in enumerate(obj):
                lines.extend(flatten(item, f"{prefix}[{i}] "))
        else:
            lines.append(f"{prefix}{obj}")
        return lines

    content = "\n".join(flatten(data))
    return [Document(
        page_content=content,
        metadata={"source": file_path, "data_type": "json"},
    )]


def load_xml(file_path: str) -> List[Document]:
    tree = ET.parse(file_path)
    root = tree.getroot()

    def elem_to_text(elem, depth=0) -> str:
        indent = "  " * depth
        parts = [f"{indent}<{elem.tag}"]
        if elem.attrib:
            attrs = " ".join(f'{k}="{v}"' for k, v in elem.attrib.items())
            parts[0] += f" {attrs}"
        parts[0] += ">"
        if elem.text and elem.text.strip():
            parts.append(f"{indent}  {elem.text.strip()}")
        for child in elem:
            parts.append(elem_to_text(child, depth + 1))
        parts.append(f"{indent}</{elem.tag}>")
        return "\n".join(parts)

    content = elem_to_text(root)
    return [Document(
        page_content=content,
        metadata={"source": file_path, "data_type": "xml"},
    )]


# ─────────────────────────────────────────────────────────────────────────────
# Structured
# ─────────────────────────────────────────────────────────────────────────────

def load_csv(file_path: str) -> List[Document]:
    docs = []
    with open(file_path, newline="", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        rows = list(reader)
        fieldnames = reader.fieldnames or []

    schema_doc = Document(
        page_content=(
            f"CSV File: {os.path.basename(file_path)}\n"
            f"Columns: {', '.join(fieldnames)}\n"
            f"Total Rows: {len(rows)}\n\n"
            f"Sample rows:\n" +
            "\n".join(str(dict(row)) for row in rows[:5])
        ),
        metadata={"source": file_path, "data_type": "csv_schema"},
    )
    docs.append(schema_doc)

    for i, row in enumerate(rows):
        content = "\n".join(f"{k}: {v}" for k, v in row.items())
        docs.append(Document(
            page_content=content,
            metadata={"source": file_path, "data_type": "csv_row", "row": i + 1},
        ))
    return docs


def load_sql_database(db_path: str) -> List[Document]:
    conn = sqlite3.connect(db_path)
    cursor = conn.cursor()
    docs: List[Document] = []

    cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
    tables = cursor.fetchall()

    for (table_name,) in tables:
        cursor.execute(f"PRAGMA table_info({table_name});")
        columns = [col[1] for col in cursor.fetchall()]

        cursor.execute(f"SELECT * FROM {table_name}")
        rows = cursor.fetchall()

        content = (
            f"Table: {table_name}\n"
            f"Columns: {', '.join(columns)}\n"
            f"Total Records: {len(rows)}\n\n"
            "Sample Records:\n"
        )
        for row in rows[:5]:
            content += f"{dict(zip(columns, row))}\n"

        docs.append(Document(
            page_content=content,
            metadata={
                "source": db_path,
                "table_name": table_name,
                "num_records": len(rows),
                "data_type": "sql_table",
            },
        ))

    conn.close()
    return docs


# ─────────────────────────────────────────────────────────────────────────────
# Auto-detect & unified entry point
# ─────────────────────────────────────────────────────────────────────────────

EXT_MAP = {
    ".txt":    "text",
    ".md":     "text",
    ".pdf":    "pdf",
    ".docx":   "docx",
    ".pptx":   "pptx",
    ".json":   "json",
    ".xml":    "xml",
    ".csv":    "csv",
    ".db":     "sql",
    ".sqlite": "sql",
}


def load_documents(source: str, source_type: str = "auto", **kwargs) -> List[Document]:
    """
    Universal document loader.

    Args:
        source:      File path, directory, or DB path.
        source_type: "auto" to detect from extension, or one of:
                     text | pdf | docx | pptx | json | xml | csv | sql | directory
        **kwargs:    Passed to the specific loader.

    Returns:
        List[Document]
    """
    if source_type == "auto":
        if os.path.isdir(source):
            source_type = "directory"
        else:
            ext = os.path.splitext(source)[1].lower()
            source_type = EXT_MAP.get(ext, "text")

    loaders = {
        "text":      load_text_file,
        "directory": load_text_directory,
        "pdf":       load_pdf,
        "docx":      load_docx,
        "pptx":      load_pptx,
        "json":      load_json,
        "xml":       load_xml,
        "csv":       load_csv,
        "sql":       load_sql_database,
    }

    if source_type not in loaders:
        raise ValueError(
            f"Unknown source_type '{source_type}'. "
            f"Valid options: {list(loaders)}"
        )

    docs = loaders[source_type](source, **kwargs)
    print(f"[parser] Loaded {len(docs)} document(s) from '{source}' "
          f"(type={source_type})")
    return docs
